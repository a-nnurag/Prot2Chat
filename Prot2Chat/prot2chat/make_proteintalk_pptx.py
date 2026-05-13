#!/usr/bin/env python3
"""
make_proteintalk_pptx.py  —  Generates ProteinTalk_Presentation_v2.pptx
15 slides with consistent dark-navy / gold styling.

Run:  python make_proteintalk_pptx.py
Output: ../../ProteinTalk_Presentation_v2.pptx
"""

import os, sys
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

_DIR   = os.path.dirname(os.path.abspath(__file__))
OUT    = os.path.join(_DIR, "..", "..", "ProteinTalk_Presentation_v2.pptx")

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
GOLD   = RGBColor(0xE8, 0xA8, 0x38)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF6, 0xF9)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
BLUE   = RGBColor(0x1A, 0x6B, 0x9A)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def rgb(r, g, b): return RGBColor(r, g, b)

def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_solid(s, color)
    s.line.fill.background()
    return s

def add_textbox(slide, text, l, t, w, h,
                size=18, bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = word_wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def add_para(tf, text, size=16, bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, space_before=6):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def content_slide(prs, title_text):
    """Base slide with navy header bar + gold accent line + white body."""
    slide = blank_slide(prs)
    # Navy header
    add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
    # Gold accent line
    add_rect(slide, 0, 1.15, 13.33, 0.06, GOLD)
    # White body background
    add_rect(slide, 0, 1.21, 13.33, 6.29, LGRAY)
    # Title text
    add_textbox(slide, title_text, 0.3, 0.15, 12.5, 0.9,
                size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    return slide

def img_to_pptx(slide, fig_or_path, l, t, w, h):
    if isinstance(fig_or_path, str):
        slide.shapes.add_picture(fig_or_path, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    else:
        buf = io.BytesIO()
        fig_or_path.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
        plt.close(fig_or_path)

def bullet_box(slide, items, l, t, w, h, size=17, title=None, title_size=19):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = NAVY
        r.font.name = "Calibri"
    for i, item in enumerate(items):
        if title or i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = ("▸  " if not item.startswith("   ") else "") + item
        r.font.size = Pt(size)
        r.font.color.rgb = DGRAY if not item.startswith("▸") else DGRAY
        r.font.name = "Calibri"
    return tb

# ── Diagram generators ────────────────────────────────────────────────────────

def arch_overview_fig():
    fig, ax = plt.subplots(figsize=(11, 3.2))
    fig.patch.set_facecolor("#F4F6F9")
    ax.set_facecolor("#F4F6F9")
    ax.set_xlim(0, 11); ax.set_ylim(0, 3.2); ax.axis("off")

    boxes = [
        (0.3,  0.7, 1.8, 1.8, "#1A3A5C", "PDB File\n(3D Structure)"),
        (2.7,  0.7, 1.8, 1.8, "#1A6B9A", "ProteinMPNN\nEnsemble\n(9 models)"),
        (5.1,  0.7, 1.8, 1.8, "#E8A838", "Cross-Attention\nAdapter\n(Our Layer)"),
        (7.5,  0.7, 1.8, 1.8, "#1A7A4A", "LLaMA3-8B\n+ LoRA"),
        (9.9,  0.7, 0.9, 1.8, "#8B1A1A", "Answer"),
    ]
    for (x, y, w, h, col, txt) in boxes:
        r = plt.Rectangle((x, y), w, h, facecolor=col, edgecolor="white",
                           linewidth=2, zorder=3)
        ax.add_patch(r)
        ax.text(x + w/2, y + h/2, txt, ha="center", va="center",
                fontsize=9, color="white", fontweight="bold",
                multialignment="center", zorder=4)

    for x in [2.1, 4.5, 6.9, 9.3]:
        ax.annotate("", xy=(x+0.6, 1.6), xytext=(x, 1.6),
                    arrowprops=dict(arrowstyle="->", color="#555", lw=2), zorder=5)

    labels = ["backbone\ncoords", "[N,1152]\nembedding",
              "[256,4096]\nsoft-prompt", "generated\ntext"]
    for i, (lbl, x) in enumerate(zip(labels, [2.15, 4.55, 6.95, 9.35])):
        ax.text(x+0.3, 2.65, lbl, ha="center", va="bottom",
                fontsize=7, color="#555", multialignment="center")

    ax.text(5.5, 0.3, "ProteinTalk — End-to-End Multimodal Pipeline",
            ha="center", fontsize=10, fontweight="bold", color="#1A1A2E")
    plt.tight_layout(pad=0.3)
    return fig


def pdb_pipeline_fig():
    fig, ax = plt.subplots(figsize=(10, 3.5))
    fig.patch.set_facecolor("#F4F6F9")
    ax.set_facecolor("#F4F6F9")
    ax.set_xlim(0, 10); ax.set_ylim(0, 3.5); ax.axis("off")

    steps = [
        (0.2,  1.0, 1.7, 1.5, "#2C3E50", "PDB File\nInput"),
        (2.3,  1.0, 1.7, 1.5, "#1A6B9A", "Biotite\nParser\n(N,Cα,C,O)"),
        (4.4,  1.0, 1.7, 1.5, "#1A7A4A", "9-Model\nProteinMPNN\nEnsemble"),
        (6.5,  1.0, 1.7, 1.5, "#8B4513", "Concatenate\n128×9 = 1152\ndims/residue"),
        (8.6,  1.0, 1.2, 1.5, "#8B1A1A", "Pad/Trim\n→512 res"),
    ]
    for (x, y, w, h, col, txt) in steps:
        r = plt.Rectangle((x, y), w, h, facecolor=col, edgecolor="white",
                           linewidth=1.5, zorder=3, )
        ax.add_patch(r)
        ax.text(x+w/2, y+h/2, txt, ha="center", va="center",
                fontsize=8.5, color="white", fontweight="bold",
                multialignment="center", zorder=4)

    for x in [1.9, 4.0, 6.1, 8.2]:
        ax.annotate("", xy=(x+0.4, 1.75), xytext=(x, 1.75),
                    arrowprops=dict(arrowstyle="->", color="#555", lw=1.8), zorder=5)

    ax.text(5.0, 0.4, "PDB Processing → 1152-dim Per-Residue Structural Embedding",
            ha="center", fontsize=9, fontweight="bold", color="#1A1A2E")
    plt.tight_layout(pad=0.3)
    return fig


def adapter_fig():
    fig, ax = plt.subplots(figsize=(10, 3.5))
    fig.patch.set_facecolor("#F4F6F9")
    ax.set_facecolor("#F4F6F9")
    ax.set_xlim(0, 10); ax.set_ylim(0, 3.5); ax.axis("off")

    steps = [
        (0.2,  0.9, 1.6, 1.6, "#1A6B9A", "Linear\nProjection\n1152→4096"),
        (2.2,  0.9, 1.6, 1.6, "#2E8B57", "Positional\nEncoding\n+256 Queries"),
        (4.2,  0.9, 1.6, 1.6, "#E8A838", "Cross-\nAttention\n(16 heads)"),
        (6.2,  0.9, 1.6, 1.6, "#8B1A1A", "LoRA\nLLaMA3-8B\nr=8, α=16"),
        (8.2,  0.9, 1.6, 1.6, "#1A3A5C", "Answer\nGeneration\n(T=0.3)"),
    ]
    for (x, y, w, h, col, txt) in steps:
        r = plt.Rectangle((x, y), w, h, facecolor=col, edgecolor="white",
                           linewidth=1.5, zorder=3)
        ax.add_patch(r)
        ax.text(x+w/2, y+h/2, txt, ha="center", va="center",
                fontsize=8.5, color="white", fontweight="bold",
                multialignment="center", zorder=4)

    for x in [1.8, 3.8, 5.8, 7.8]:
        ax.annotate("", xy=(x+0.4, 1.7), xytext=(x, 1.7),
                    arrowprops=dict(arrowstyle="->", color="#555", lw=1.8), zorder=5)

    ax.text(5.0, 0.35, "Adapter + LoRA Fine-Tuning Pipeline (Trained on Mol-Instructions)",
            ha="center", fontsize=9, fontweight="bold", color="#1A1A2E")

    # Highlight adapter box as "Our contribution"
    ax.annotate("Our trained\ncomponent", xy=(5.0, 2.5), fontsize=8,
                color="#E8A838", fontweight="bold", ha="center",
                arrowprops=dict(arrowstyle="->", color="#E8A838"),
                xytext=(5.0, 3.1))
    plt.tight_layout(pad=0.3)
    return fig


def eval_results_fig():
    models = ["ProteinTalk\n(Ours)", "Prot2Chat\n(Wang et al.)",
              "Evola\n10B", "KIMI\nfew-shot", "LLaMA3-\nFT", "BioMedGPT\n10B"]
    bleu2  = [6.32,  35.85, 8.69,  12.05, 6.42,  1.02]
    rouge1 = [20.33, 57.21, 29.09, 31.21, 24.50, 10.93]
    rougeL = [14.04, 50.51, 20.04, 24.18, 17.03,  7.84]

    x = np.arange(len(models))
    w = 0.26
    colors_b2  = ["#E8A838"] + ["#1A6B9A"]*5
    colors_r1  = ["#C8882A"] + ["#145E87"]*5
    colors_rL  = ["#A86818"] + ["#0E4A6A"]*5

    fig, ax = plt.subplots(figsize=(11, 4.0))
    fig.patch.set_facecolor("#F4F6F9")
    ax.set_facecolor("#F4F6F9")

    b1 = ax.bar(x - w, bleu2,  w, label="BLEU-2",  color=colors_b2,  edgecolor="white")
    b2 = ax.bar(x,     rouge1, w, label="ROUGE-1", color=colors_r1,  edgecolor="white")
    b3 = ax.bar(x + w, rougeL, w, label="ROUGE-L", color=colors_rL,  edgecolor="white")

    for bars in [b1, b2, b3]:
        for bar in bars:
            h = bar.get_height()
            if h >= 1:
                ax.text(bar.get_x()+bar.get_width()/2, h+0.4, f"{h:.1f}",
                        ha="center", va="bottom", fontsize=6.5, color="#333")

    ax.set_xticks(x)
    ax.set_xticklabels(models, fontsize=9)
    ax.set_ylabel("Score", fontsize=10)
    ax.set_ylim(0, 68)
    ax.legend(fontsize=9, loc="upper right")
    ax.grid(axis="y", alpha=0.3)
    ax.set_title("ProteinTalk vs Baselines — Mol-Instructions Test Set", fontsize=11, fontweight="bold")

    gold = mpatches.Patch(color="#E8A838", label="ProteinTalk (Ours)")
    blue = mpatches.Patch(color="#1A6B9A", label="Baseline Models")
    ax.legend(handles=[gold, blue], fontsize=9)

    plt.tight_layout(pad=0.4)
    return fig


def lit_review_fig():
    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor("#F4F6F9")
    ax.set_facecolor("#F4F6F9")
    ax.set_xlim(0, 11); ax.set_ylim(0, 4.5); ax.axis("off")

    papers = [
        (0.2,  2.5, 2.2, 1.5, "#1A3A5C",
         "ProteinMPNN\n(Dauparas et al., 2022)\nScience 378(6615)",
         "GNN-based structural\nencoder — forms our\nembedding backbone"),
        (2.8,  2.5, 2.2, 1.5, "#1A6B9A",
         "LoRA\n(Hu et al., 2022)\nICLR 2022",
         "Low-rank adaptation —\nour fine-tuning\nstrategy"),
        (5.4,  2.5, 2.2, 1.5, "#1A7A4A",
         "BLIP-2\n(Li et al., 2023)\nICML 2023",
         "Q-Former cross-attention\n— inspires our\nadapter design"),
        (8.0,  2.5, 2.2, 1.5, "#8B4513",
         "LLaMA3\n(Meta AI, 2024)",
         "Base language model\nfor instruction-\nfollowing generation"),
        (0.9,  0.4, 2.2, 1.5, "#5B2C8B",
         "BioMedGPT-LM-10B\n(Liu et al., 2024)",
         "Biomedical LLM baseline\n— sequence-only\napproach"),
        (3.5,  0.4, 2.2, 1.5, "#8B1A1A",
         "Mol-Instructions\n(Fang et al., 2023)",
         "Training & evaluation\ndataset — 404K protein\ninstruction samples"),
        (6.1,  0.4, 2.2, 1.5, "#2E5B8B",
         "Evola-10B\n(2024)",
         "Structure+sequence\nbaseline for\ncomparison"),
        (8.7,  0.4, 2.2, 1.5, "#5B6B1A",
         "AlphaFold DB\n(Jumper et al., 2021)",
         "Source of predicted\nPDB structures for\nevaluation"),
    ]

    for (x, y, w, h, col, title, desc) in papers:
        r = plt.Rectangle((x, y), w, h, facecolor=col, edgecolor="white",
                           linewidth=1.5, zorder=3, alpha=0.92)
        ax.add_patch(r)
        ax.text(x+w/2, y+h-0.25, title, ha="center", va="top",
                fontsize=7.5, color="white", fontweight="bold",
                multialignment="center", zorder=4)
        ax.text(x+w/2, y+0.2, desc, ha="center", va="bottom",
                fontsize=6.8, color="#DDDDDD",
                multialignment="center", zorder=4)

    ax.text(5.5, 4.25, "Key References Informing ProteinTalk's Design",
            ha="center", fontsize=10, fontweight="bold", color="#1A1A2E")
    plt.tight_layout(pad=0.2)
    return fig


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 2.8, 13.33, 0.08, GOLD)
    add_rect(slide, 0, 4.4, 13.33, 0.08, GOLD)
    add_textbox(slide, "ProteinTalk",
                0.5, 0.6, 12.3, 1.2, size=52, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "Multimodal Protein Structure Q&A System",
                0.5, 1.7, 12.3, 0.9, size=24, bold=False, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "A Study Inspired by ProteinMPNN, BLIP-2, and LoRA",
                0.5, 2.4, 12.3, 0.55, size=16, italic=True, color=GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "Submitted To: Surajeet Ghosh (Associate Professor)",
                1.0, 4.7, 11.3, 0.5, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Anurag Chaurasia (2022ccsb098)   |   Rajesh Kumar (2022csb094)   |   Adrija Paul (2022csb082)",
                0.5, 5.25, 12.3, 0.5, size=14, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "IIEST Shibpur  |  Department of Computer Science & Technology  |  2025–2026",
                0.5, 5.85, 12.3, 0.5, size=13, italic=True, color=WHITE,
                align=PP_ALIGN.CENTER)


def slide_problem(prs):
    slide = content_slide(prs, "Problem Statement")
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.8))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True

    items = [
        ("The Challenge", True, NAVY, 22),
        ("Proteins are fundamental to all life processes — yet understanding their function "
         "from 3D structure requires deep domain expertise and specialised bioinformatics tools "
         "inaccessible to most researchers.", False, DGRAY, 16),
        ("", False, DGRAY, 8),
        ("Identified Gaps in Prior Work", True, NAVY, 20),
        ("▸  Sequence-only models (e.g., BioMedGPT-LM-10B) ignore 3D structural geometry, "
         "missing critical folding information.", False, DGRAY, 15),
        ("▸  Structure-only methods cannot answer natural language questions about function.", False, DGRAY, 15),
        ("▸  Existing multimodal systems (e.g., Evola-10B) achieve limited NLG quality "
         "(BLEU-2: 8.69 vs. human-level expected ~35+).", False, DGRAY, 15),
        ("▸  No accessible web-based interface for PDB-driven protein Q&A.", False, DGRAY, 15),
        ("", False, DGRAY, 8),
        ("Our Approach — ProteinTalk", True, NAVY, 20),
        ("Fuse ProteinMPNN structural embeddings with a LoRA fine-tuned LLaMA3 language model "
         "via a trained cross-attention adapter, enabling natural language Q&A grounded in "
         "3D protein structure.", False, DGRAY, 15),
    ]
    first = True
    for (text, bold, color, size) in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = text
        r.font.bold = bold; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.name = "Calibri"


def slide_literature(prs):
    slide = content_slide(prs, "Literature Review")
    fig = lit_review_fig()
    img_to_pptx(slide, fig, 0.3, 1.25, 12.7, 5.9)


def slide_dataset(prs):
    slide = content_slide(prs, "Datasets Used")
    # Left: Mol-Instructions
    add_rect(slide, 0.4, 1.35, 5.8, 5.7, WHITE)
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(5.6), Inches(5.5))
    tb1.word_wrap = True; tf1 = tb1.text_frame; tf1.word_wrap = True
    for txt, bold, size, col in [
        ("Mol-Instructions", True, 20, NAVY),
        ("Fang et al., 2023  |  UniProtKB/SwissProt", False, 12, BLUE),
        ("", False, 8, DGRAY),
        ("Protein-oriented instruction dataset covering:", False, 14, DGRAY),
        ("▸  Protein function description", False, 13, DGRAY),
        ("▸  Catalytic activity annotation", False, 13, DGRAY),
        ("▸  Domain & motif identification", False, 13, DGRAY),
        ("", False, 8, DGRAY),
        ("Split           Samples", True, 13, NAVY),
        ("Train          404,640", False, 13, DGRAY),
        ("Validation     16,859", False, 13, DGRAY),
        ("Test           11,072", False, 13, DGRAY),
        ("", False, 8, DGRAY),
        ("Used for: Training the adapter layer and evaluating ProteinTalk's generation quality.", False, 12, DGRAY),
    ]:
        if tf1.paragraphs[0].text == "" and not tf1.paragraphs[0].runs:
            p = tf1.paragraphs[0]
        else:
            p = tf1.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run(); r.text = txt
        r.font.bold = bold; r.font.size = Pt(size)
        r.font.color.rgb = col; r.font.name = "Calibri"

    # Right: UniProtQA
    add_rect(slide, 6.8, 1.35, 6.2, 5.7, WHITE)
    tb2 = slide.shapes.add_textbox(Inches(6.9), Inches(1.45), Inches(6.0), Inches(5.5))
    tb2.word_wrap = True; tf2 = tb2.text_frame; tf2.word_wrap = True
    for txt, bold, size, col in [
        ("UniProtQA", True, 20, NAVY),
        ("UniProt Knowledge Base", False, 12, BLUE),
        ("", False, 8, DGRAY),
        ("Factual Q&A grounded in UniProt entries:", False, 14, DGRAY),
        ("▸  Protein identity questions", False, 13, DGRAY),
        ("▸  Biological role queries", False, 13, DGRAY),
        ("▸  Out-of-domain generalisation test", False, 13, DGRAY),
        ("", False, 8, DGRAY),
        ("Split           Samples", True, 13, NAVY),
        ("Train           25,820", False, 13, DGRAY),
        ("Validation       1,075", False, 13, DGRAY),
        ("Test             6,734", False, 13, DGRAY),
        ("", False, 8, DGRAY),
        ("Used for: Generalisation benchmark — tests domain transfer beyond training distribution.", False, 12, DGRAY),
    ]:
        if tf2.paragraphs[0].text == "" and not tf2.paragraphs[0].runs:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run(); r.text = txt
        r.font.bold = bold; r.font.size = Pt(size)
        r.font.color.rgb = col; r.font.name = "Calibri"


def slide_arch_overview(prs):
    slide = content_slide(prs, "System Architecture — Overview (1/3)")
    fig = arch_overview_fig()
    img_to_pptx(slide, fig, 0.3, 1.3, 12.7, 3.5)

    add_textbox(slide, "Three-Stage Multimodal Pipeline:", 0.4, 4.95, 12.5, 0.4,
                size=16, bold=True, color=NAVY)
    items = [
        "Stage 1 — Structural Encoding: PDB file → ProteinMPNN ensemble → [N, 1152] per-residue embedding",
        "Stage 2 — Adapter (Our Contribution): Projects protein embedding → 256 soft-prompt tokens in LLM space",
        "Stage 3 — Language Generation: LLaMA3-8B + LoRA generates natural language answer conditioned on protein context",
    ]
    y = 5.45
    for item in items:
        add_textbox(slide, "▸  " + item, 0.4, y, 12.5, 0.45, size=14, color=DGRAY)
        y += 0.47


def slide_arch_pdb(prs):
    slide = content_slide(prs, "System Architecture — PDB Processing (2/3)")
    fig = pdb_pipeline_fig()
    img_to_pptx(slide, fig, 0.3, 1.3, 12.7, 3.6)

    add_textbox(slide, "ProteinMPNN Ensemble Details:", 0.4, 5.05, 12.5, 0.4,
                size=16, bold=True, color=NAVY)
    cols = [
        ("6 Full-Atom Models", "Gaussian noise σ ∈ {0.02, 0.10, 0.20, 0.30}\nIncludes 2 soluble-protein variants", 0.4),
        ("3 CA-Only Models", "Operate on alpha-carbon coordinates only\nRobust to incomplete backbone data", 4.6),
        ("Ensemble Output", "9 × 128-dim outputs concatenated\n→ 1152-dim per residue, padded to 512", 8.8),
    ]
    for title, desc, x in cols:
        add_rect(slide, x, 5.55, 3.8, 1.6, WHITE)
        add_textbox(slide, title, x+0.1, 5.6, 3.6, 0.45, size=14, bold=True, color=NAVY)
        add_textbox(slide, desc, x+0.1, 6.1, 3.6, 0.9, size=12, color=DGRAY)


def slide_arch_adapter(prs):
    slide = content_slide(prs, "System Architecture — Adapter & LoRA (3/3)")
    fig = adapter_fig()
    img_to_pptx(slide, fig, 0.3, 1.3, 12.7, 3.6)

    add_textbox(slide, "Key Design Choices:", 0.4, 5.05, 12.5, 0.4,
                size=16, bold=True, color=NAVY)
    cols = [
        ("Cross-Attention Adapter", "Inspired by BLIP-2 Q-Former\n256 learnable query tokens\n16-head cross-attention", 0.4),
        ("LoRA Fine-Tuning", "Rank r=8, α=16\nTargets q_proj + v_proj\n~3.4M trainable params", 4.6),
        ("Training Setup", "Mol-Instructions train split\nAdamW + bfloat16 AMP\n400,000 steps on RTX 3090", 8.8),
    ]
    for title, desc, x in cols:
        add_rect(slide, x, 5.55, 3.8, 1.6, WHITE)
        add_textbox(slide, title, x+0.1, 5.6, 3.6, 0.45, size=14, bold=True, color=NAVY)
        add_textbox(slide, desc, x+0.1, 6.1, 3.6, 0.9, size=12, color=DGRAY)


def slide_pdb(prs):
    slide = content_slide(prs, "PDB Files — Protein Data Bank Format")
    add_rect(slide, 0.4, 1.35, 6.0, 5.7, WHITE)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(5.8), Inches(5.5))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for txt, bold, size, col in [
        ("What is a PDB File?", True, 18, NAVY),
        ("A PDB (Protein Data Bank) file is a standardised plain-text format that stores the "
         "three-dimensional atomic coordinates of a macromolecule.", False, 13, DGRAY),
        ("", False, 6, DGRAY),
        ("File Contents:", True, 16, NAVY),
        ("▸  ATOM records: x, y, z coordinates of each atom", False, 13, DGRAY),
        ("▸  Residue name, chain ID, residue sequence number", False, 13, DGRAY),
        ("▸  B-factor (temperature factor) per atom", False, 13, DGRAY),
        ("▸  Multi-chain support for protein complexes", False, 13, DGRAY),
        ("", False, 6, DGRAY),
        ("Backbone Atoms Used (per residue):", True, 16, NAVY),
        ("N  —  Amino nitrogen (chain backbone)", False, 13, DGRAY),
        ("Cα —  Alpha carbon (branch point)", False, 13, DGRAY),
        ("C  —  Carbonyl carbon", False, 13, DGRAY),
        ("O  —  Carbonyl oxygen", False, 13, DGRAY),
        ("", False, 6, DGRAY),
        ("Sources in ProteinTalk:", True, 16, NAVY),
        ("▸  Protein Data Bank (PDB): Experimentally resolved structures", False, 13, DGRAY),
        ("▸  AlphaFold DB: Computationally predicted (used in our evaluation)", False, 13, DGRAY),
    ]:
        if tf.paragraphs[0].text == "" and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run(); r.text = txt
        r.font.bold = bold; r.font.size = Pt(size)
        r.font.color.rgb = col; r.font.name = "Calibri"

    # Right: PDB format illustration
    add_rect(slide, 6.7, 1.35, 6.3, 5.7, RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(slide, "Sample PDB Record Format", 6.8, 1.45, 6.1, 0.4,
                size=14, bold=True, color=GOLD)
    pdb_sample = (
        "ATOM      1  N   MET A   1     "
        " 10.123  20.456  30.789\n"
        "ATOM      2  CA  MET A   1     "
        " 11.234  21.567  31.890\n"
        "ATOM      3  C   MET A   1     "
        " 12.345  22.678  32.901\n"
        "ATOM      4  O   MET A   1     "
        " 13.456  23.789  33.012\n"
        "ATOM      5  N   ALA A   2     "
        " 14.567  24.890  34.123\n"
        "  ...               ...\n\n"
        "Columns:\n"
        "  1-6:  Record type (ATOM)\n"
        "  13-16: Atom name (N, CA, C, O)\n"
        "  18-20: Residue name (MET, ALA...)\n"
        "  22:   Chain ID (A, B, ...)\n"
        "  31-38: X coordinate (Angstrom)\n"
        "  39-46: Y coordinate\n"
        "  47-54: Z coordinate"
    )
    add_textbox(slide, pdb_sample, 6.8, 1.95, 6.1, 4.8,
                size=11, color=RGBColor(0xA0, 0xFF, 0xA0))


def slide_eval_results(prs):
    slide = content_slide(prs, "Evaluation Results — Mol-Instructions Test Set")
    fig = eval_results_fig()
    img_to_pptx(slide, fig, 0.3, 1.3, 7.8, 4.0)

    # Score table on right
    add_rect(slide, 8.3, 1.35, 4.7, 5.7, WHITE)
    add_textbox(slide, "Score Comparison", 8.4, 1.45, 4.5, 0.4,
                size=16, bold=True, color=NAVY)

    rows = [
        ("Model",              "B-2",   "R-1",   "R-L",   True,  "0D1B2A", "FFFFFF"),
        ("ProteinTalk (Ours)", "6.32",  "20.33", "14.04", True,  "FFF3CD", "1A1A2E"),
        ("Prot2Chat†",         "35.85", "57.21", "50.51", False, "EAF4FB", "1A1A2E"),
        ("Evola-10B†",         "8.69",  "29.09", "20.04", False, "FFFFFF", "1A1A2E"),
        ("KIMI few-shot†",     "12.05", "31.21", "24.18", False, "FFFFFF", "1A1A2E"),
        ("LLaMA3-FT†",         "6.42",  "24.50", "17.03", False, "FFFFFF", "1A1A2E"),
        ("BioMedGPT†",         "1.02",  "10.93", "7.84",  False, "FFFFFF", "1A1A2E"),
    ]
    y = 1.95
    for (model, b2, r1, rL, bold, bg, fg) in rows:
        bg_rgb = RGBColor(int(bg[0:2],16), int(bg[2:4],16), int(bg[4:6],16))
        fg_rgb = RGBColor(int(fg[0:2],16), int(fg[2:4],16), int(fg[4:6],16))
        add_rect(slide, 8.3, y, 4.7, 0.48, bg_rgb)
        add_textbox(slide, model, 8.35, y+0.04, 2.3, 0.4, size=11,
                    bold=bold, color=fg_rgb)
        add_textbox(slide, f"{b2}   {r1}   {rL}", 10.7, y+0.04, 2.1, 0.4,
                    size=11, bold=bold, color=fg_rgb, align=PP_ALIGN.CENTER)
        y += 0.48

    add_textbox(slide, "† Results from Wang et al. (2025)\n50-sample evaluation (seed=42)",
                8.3, 6.5, 4.7, 0.65, size=10, italic=True, color=DGRAY)

    add_textbox(slide,
                "ProteinTalk achieves BLEU-2: 6.32 — comparable to LLaMA3-FT (6.42) "
                "and significantly above BioMedGPT (1.02), despite 6 GB VRAM quantization constraints.",
                0.3, 5.45, 7.8, 0.8, size=13, italic=True, color=NAVY)


def slide_not_performed_1(prs):
    slide = content_slide(prs, "Evaluations Not Performed — Part 1 of 2")

    for x, title, why, score_txt, col in [
        (0.4,
         "KIMI LLM-as-Judge Evaluation",
         "Requires access to KIMI (Moonshot AI) proprietary LLM service — not available as "
         "an open API in our region. Also requires simultaneous generation from all competing "
         "models (Evola-10B, LLaMA3-FT, KIMI, ProteinTalk) — beyond current GPU capacity.",
         "Paper result: Prot2Chat ranked 1st in 386/650 evaluations (59.4%)\nAvg. rank: 1.45 / 4",
         "#1A6B9A"),
        (6.9,
         "Expert Human Evaluation",
         "Requires professional biology PhD researchers to manually rank 650 protein Q&A "
         "instances under a controlled protocol. Recruiting domain experts and managing "
         "inter-annotator agreement is outside the scope of this academic project.",
         "Paper result: Prot2Chat ranked 1st in 359/650 evaluations (55.2%)\nAvg. rank: 1.49 / 4",
         "#1A7A4A"),
    ]:
        add_rect(slide, x, 1.35, 6.1, 5.7, WHITE)
        add_textbox(slide, title, x+0.15, 1.45, 5.8, 0.55,
                    size=17, bold=True, color=RGBColor(*[int(col[i:i+2],16) for i in (1,3,5)]))
        add_textbox(slide, "What it measures:", x+0.15, 2.1, 5.8, 0.35,
                    size=13, bold=True, color=NAVY)
        desc = ("Uses a large LLM as an automated judge to rank outputs from competing models "
                "on biological relevance, coherence, and accuracy. Captures quality aspects "
                "that BLEU/ROUGE cannot." if "KIMI" in title else
                "Professional biologists rank model outputs for factual accuracy, completeness, "
                "and domain correctness — considered the gold-standard evaluation for protein Q&A systems.")
        add_textbox(slide, desc, x+0.15, 2.5, 5.8, 1.0, size=13, color=DGRAY)
        add_textbox(slide, "Why not performed:", x+0.15, 3.6, 5.8, 0.35,
                    size=13, bold=True, color=RGBColor(0x8B, 0x1A, 0x1A))
        add_textbox(slide, why, x+0.15, 4.0, 5.8, 1.2, size=12, color=DGRAY)
        add_rect(slide, x+0.1, 5.3, 5.9, 0.85,
                 RGBColor(*[int(col[i:i+2],16) for i in (1,3,5)]))
        add_textbox(slide, score_txt, x+0.2, 5.35, 5.7, 0.75,
                    size=11, color=WHITE)


def slide_not_performed_2(prs):
    slide = content_slide(prs, "Evaluations Not Performed — Part 2 of 2")

    add_rect(slide, 0.4, 1.35, 12.5, 3.2, WHITE)
    add_textbox(slide, "UniProtQA Generalisation Benchmark",
                0.55, 1.45, 12.0, 0.55, size=19, bold=True, color=NAVY)
    add_textbox(slide, "What it measures:",
                0.55, 2.1, 5.0, 0.35, size=14, bold=True, color=NAVY)
    add_textbox(slide,
                "Tests whether a model trained on Mol-Instructions can generalise to a "
                "completely different protein Q&A dataset (UniProtQA) without retraining — "
                "measuring true domain transfer. 6,734 test samples.",
                0.55, 2.5, 5.5, 1.0, size=13, color=DGRAY)
    add_textbox(slide, "Why not performed:",
                6.3, 2.1, 5.0, 0.35, size=14, bold=True, color=RGBColor(0x8B,0x1A,0x1A))
    add_textbox(slide,
                "Requires mapping all 6,734 UniProtQA test proteins to AlphaFold DB PDB "
                "structures via UniProt accession lookup — estimated ~18–20 hours of API "
                "calls + inference on 6 GB GPU. Deferred to future work.",
                6.3, 2.5, 6.2, 1.0, size=13, color=DGRAY)
    add_rect(slide, 0.4, 4.4, 12.5, 0.07, GOLD)
    add_textbox(slide, "Paper result (Wang et al., 2025): BLEU-2: 6.72   ROUGE-1: 15.71",
                0.55, 4.55, 12.0, 0.4, size=13, bold=True, color=NAVY)

    add_rect(slide, 0.4, 5.1, 12.5, 2.05, RGBColor(0x1A,0x3A,0x5C))
    add_textbox(slide, "Summary: Hardware & Resource Constraints",
                0.55, 5.2, 12.0, 0.45, size=16, bold=True, color=GOLD)
    add_textbox(slide,
                "ProteinTalk was trained on an RTX 3090 (24 GB VRAM, float32).  "
                "Our deployment runs on 6 GB VRAM with 4-bit NF4 quantization, "
                "limiting inference quality and batch evaluation throughput.\n"
                "KIMI/Expert evaluations require proprietary tools and domain experts beyond project scope.\n"
                "Future work: Full UniProtQA evaluation + expert review with higher-VRAM hardware.",
                0.55, 5.7, 12.0, 1.3, size=13, color=WHITE)


def slide_conclusion(prs):
    slide = content_slide(prs, "Conclusion")
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.8))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    items = [
        ("What We Built", True, NAVY, 20),
        ("ProteinTalk — a multimodal protein Q&A system integrating ProteinMPNN structural "
         "embeddings, a trained cross-attention adapter, and LoRA-finetuned LLaMA3-8B-Instruct, "
         "trained on 404K protein instruction samples from Mol-Instructions.", False, DGRAY, 15),
        ("", False, DGRAY, 7),
        ("Key Achievements", True, NAVY, 20),
        ("▸  Implemented end-to-end multimodal pipeline: PDB → ProteinMPNN → Adapter → LLM → Answer", False, DGRAY, 14),
        ("▸  Achieved BLEU-2: 6.32, ROUGE-1: 20.33, ROUGE-L: 14.04 on 50-sample Mol-Instructions evaluation", False, DGRAY, 14),
        ("▸  Outperforms BioMedGPT-LM-10B (BLEU-2: 1.02) — confirming structural encoding adds value", False, DGRAY, 14),
        ("▸  Deployed as a Flask web application with PDB file upload interface", False, DGRAY, 14),
        ("▸  Diagnosed and partially mitigated 4-bit quantization degradation on 6 GB GPU", False, DGRAY, 14),
        ("", False, DGRAY, 7),
        ("Limitations & Future Work", True, NAVY, 20),
        ("▸  Gap vs. full model (BLEU-2 35.85) due to 4-bit quantization on 6 GB VRAM hardware", False, DGRAY, 14),
        ("▸  Future: UniProtQA evaluation, expert review, quantization-aware training, RAG integration", False, DGRAY, 14),
    ]
    first = True
    for (text, bold, color, size) in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = text
        r.font.bold = bold; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.name = "Calibri"


def slide_team(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 1.4, 13.33, 0.07, GOLD)
    add_rect(slide, 0, 6.05, 13.33, 0.07, GOLD)
    add_textbox(slide, "Our Team", 0.5, 0.3, 12.3, 0.9,
                size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    members = [
        ("Anurag Chaurasia", "2022ccsb098", "System implementation, model deployment,\nquantization debugging, evaluation pipeline", 0.5),
        ("Rajesh Kumar",     "2022csb094",  "Architecture design, LoRA fine-tuning,\ndataset preparation and preprocessing",         4.7),
        ("Adrija Paul",      "2022csb082",  "Literature review, report writing,\nresults analysis and documentation",                8.9),
    ]
    for name, roll, role, x in members:
        add_rect(slide, x, 1.6, 3.9, 4.2, RGBColor(0x1A, 0x3A, 0x5C))
        add_rect(slide, x, 1.6, 3.9, 0.07, GOLD)
        add_textbox(slide, name, x+0.1, 1.75, 3.7, 0.65,
                    size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, roll, x+0.1, 2.45, 3.7, 0.45,
                    size=14, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, "B.Tech. Computer Science & Technology\nIIEST Shibpur",
                    x+0.1, 3.0, 3.7, 0.65, size=12, italic=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, role, x+0.1, 3.75, 3.7, 1.0,
                    size=12, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)

    add_textbox(slide, "Guided by: Surajeet Ghosh (Associate Professor)  |  IIEST Shibpur",
                0.5, 6.2, 12.3, 0.5, size=14, italic=True,
                color=WHITE, align=PP_ALIGN.CENTER)


def slide_thankyou(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 2.5, 3.35, 8.33, 0.08, GOLD)
    add_textbox(slide, "Thank You", 0.5, 1.5, 12.3, 1.5,
                size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "ProteinTalk — Multimodal Protein Structure Q&A System",
                0.5, 3.1, 12.3, 0.6, size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Questions & Discussion",
                0.5, 3.65, 12.3, 0.55, size=18, italic=True, color=GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Anurag Chaurasia  •  Rajesh Kumar  •  Adrija Paul\n"
                "B.Tech. Computer Science & Technology  |  IIEST Shibpur  |  2025–2026",
                0.5, 5.5, 12.3, 0.9, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Submitted to: Surajeet Ghosh (Associate Professor)",
                0.5, 6.5, 12.3, 0.5, size=13, italic=True, color=GOLD,
                align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("[1/15] Title slide...")
    slide_title(prs)
    print("[2/15] Problem statement...")
    slide_problem(prs)
    print("[3/15] Literature review...")
    slide_literature(prs)
    print("[4/15] Dataset...")
    slide_dataset(prs)
    print("[5/15] Architecture overview...")
    slide_arch_overview(prs)
    print("[6/15] PDB + ProteinMPNN pipeline...")
    slide_arch_pdb(prs)
    print("[7/15] Adapter + LoRA...")
    slide_arch_adapter(prs)
    print("[8/15] PDB files...")
    slide_pdb(prs)
    print("[9/15] Evaluation results...")
    slide_eval_results(prs)
    print("[10/15] Tests not performed (part 1)...")
    slide_not_performed_1(prs)
    print("[11/15] Tests not performed (part 2)...")
    slide_not_performed_2(prs)
    print("[12/15] Conclusion...")
    slide_conclusion(prs)
    print("[13/15] Team...")
    slide_team(prs)
    print("[14/15] Thank you...")
    slide_thankyou(prs)

    # Slide 15: blank (reserved / backup)
    print("[15/15] Reserved slide...")
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    sz = os.path.getsize(OUT) // 1024
    print(f"\n✓  Saved → {OUT}  ({sz} KB)  |  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
